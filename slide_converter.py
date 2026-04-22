#!/usr/bin/env python3
"""
NotebookLM PDF → Editable PPTX Converter
=========================================
步驟 1: 擷取文字（位置、大小）
步驟 2: 擷取前景圖案（可移動物件）
步驟 3: AI 重新生成背景（Stability AI Inpainting）

使用方式:
  python slide_converter.py input.pdf output.pptx

環境變數:
  STABILITY_API_KEY  - Stability AI API 金鑰
                       (取得: https://platform.stability.ai/account/keys)
"""

import os
import sys
import io
import time
import requests
import numpy as np
import cv2
import easyocr
from pathlib import Path
from PIL import Image, ImageDraw, ImageFilter
from pdf2image import convert_from_path
from rembg import remove as rembg_remove, new_session as rembg_session
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ─────────────────────────────────────────────────────────────────────────────
# 設定區 (Configuration)
# ─────────────────────────────────────────────────────────────────────────────

STABILITY_API_KEY = os.environ.get("STABILITY_API_KEY", "")
PDF_PATH          = "input.pdf"
OUTPUT_PPTX       = "output.pptx"
ASSETS_DIR        = Path("slide_assets")   # 中間檔案存放資料夾

DPI               = 150    # PDF 轉圖解析度（150 = 清晰；200 = 高品質但較慢）
OCR_LANGUAGES     = ["ch_tra", "en"]  # 繁中+英 (簡中改為 "ch_sim")
USE_GPU           = False  # 若有 NVIDIA GPU 可改為 True 加速 OCR
INPAINT_PADDING   = 12    # 遮罩向外擴展像素數
MASK_BLUR         = 4     # 遮罩邊緣模糊半徑
MIN_OBJ_RATIO     = 0.0005 # 前景物件最小面積比（相對於整張投影片）
MAX_OBJ_RATIO     = 0.45   # 前景物件最大面積比

# ─────────────────────────────────────────────────────────────────────────────
# 工具函數
# ─────────────────────────────────────────────────────────────────────────────

def px_scale(px: int, src_dim: int, dst_emu: int) -> int:
    """像素座標 → PowerPoint EMU 單位"""
    return int(px / src_dim * dst_emu)


def merge_boxes(boxes: list, gap: int = 20) -> list:
    """合併相鄰或重疊的邊界框"""
    if not boxes:
        return []
    changed = True
    while changed:
        changed = False
        result, used = [], [False] * len(boxes)
        for i, (ax, ay, aw, ah) in enumerate(boxes):
            if used[i]:
                continue
            mx, my, mw, mh = ax, ay, aw, ah
            for j, (bx, by, bw, bh) in enumerate(boxes):
                if i == j or used[j]:
                    continue
                if not (bx > mx + mw + gap or bx + bw + gap < mx or
                        by > my + mh + gap or by + bh + gap < my):
                    nx  = min(mx, bx)
                    ny  = min(my, by)
                    nx2 = max(mx + mw, bx + bw)
                    ny2 = max(my + mh, by + bh)
                    mx, my, mw, mh = nx, ny, nx2 - nx, ny2 - ny
                    used[j] = True
                    changed = True
            result.append((mx, my, mw, mh))
            used[i] = True
        boxes = result
    return boxes

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 0: PDF → 每頁高解析 PNG
# ─────────────────────────────────────────────────────────────────────────────

def pdf_to_images(pdf_path: str, first_page: int = None, last_page: int = None) -> list:
    print(f"[步驟 0] 將 PDF 轉換為圖片（{DPI} DPI）...")
    kwargs = {"dpi": DPI, "poppler_path": r"C:\Program Files\poppler\Library\bin"}
    if first_page: kwargs["first_page"] = first_page
    if last_page:  kwargs["last_page"]  = last_page
    pages = convert_from_path(pdf_path, **kwargs)
    print(f"        → 共 {len(pages)} 頁")
    return pages

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 1: OCR 文字偵測（保留位置與大小）
# ─────────────────────────────────────────────────────────────────────────────

def detect_text(image: Image.Image, reader: easyocr.Reader) -> list:
    """
    回傳格式:
    [{"text": str, "rect": (x, y, w, h), "conf": float, "bbox": list}, ...]
    """
    img_np = np.array(image.convert("RGB"))
    results = reader.readtext(img_np, detail=1, paragraph=False)

    text_data = []
    for bbox, text, conf in results:
        if conf < 0.3 or not text.strip():
            continue
        pts = np.array(bbox, dtype=np.int32)
        x, y, w, h = cv2.boundingRect(pts)
        text_data.append({"text": text.strip(), "rect": (x, y, w, h),
                           "conf": conf, "bbox": bbox})

    print(f"        → 偵測到 {len(text_data)} 個文字區塊")
    return text_data

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 2: 前景圖案偵測
# ─────────────────────────────────────────────────────────────────────────────

def detect_foreground_objects(image: Image.Image, text_data: list,
                              rembg_sess=None) -> tuple:
    """
    使用 rembg AI 模型偵測前景圖形元素。
    回傳:
      fg_boxes  : [(x, y, w, h), ...]  邊界框列表
      fg_alpha  : np.ndarray (H, W)    精確 alpha 遮罩（0~255）
    """
    h_img, w_img = image.height, image.width

    # ── rembg：去除背景，取得 RGBA 影像 ──────────────────────────
    rgba = rembg_remove(image.convert("RGB"), session=rembg_sess)
    alpha = np.array(rgba)[:, :, 3]      # alpha channel = 前景

    # 建立文字遮罩（文字不算前景物件）
    text_mask = np.zeros((h_img, w_img), dtype=np.uint8)
    for td in text_data:
        x, y, w, h = td["rect"]
        pad = 15
        x1, y1 = max(0, x - pad), max(0, y - pad)
        x2, y2 = min(w_img, x + w + pad), min(h_img, y + h + pad)
        text_mask[y1:y2, x1:x2] = 255

    # 從 alpha 遮罩排除文字區域
    fg_mask = alpha.copy()
    fg_mask[text_mask > 0] = 0

    # 二值化
    _, fg_bin = cv2.threshold(fg_mask, 30, 255, cv2.THRESH_BINARY)

    # 排除靠近邊緣的細小裝飾（留 3% 邊緣緩衝）
    edge_buf_x = int(w_img * 0.03)
    edge_buf_y = int(h_img * 0.03)
    fg_bin[:edge_buf_y, :]  = 0
    fg_bin[-edge_buf_y:, :] = 0
    fg_bin[:, :edge_buf_x]  = 0
    fg_bin[:, -edge_buf_x:] = 0

    # 形態學清理
    k = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (7, 7))
    fg_bin = cv2.morphologyEx(fg_bin, cv2.MORPH_CLOSE, k, iterations=3)
    fg_bin = cv2.morphologyEx(fg_bin, cv2.MORPH_OPEN,  k, iterations=2)

    # 找輪廓 → 邊界框
    contours, _ = cv2.findContours(fg_bin, cv2.RETR_EXTERNAL,
                                   cv2.CHAIN_APPROX_SIMPLE)

    total_px = h_img * w_img
    raw_boxes = []
    for cnt in contours:
        area = cv2.contourArea(cnt)
        if total_px * MIN_OBJ_RATIO < area < total_px * MAX_OBJ_RATIO:
            x, y, w, h = cv2.boundingRect(cnt)
            raw_boxes.append((x, y, w, h))

    merged = merge_boxes(raw_boxes, gap=30)
    filtered = [(x, y, w, h) for x, y, w, h in merged
                if 0.05 < w / max(h, 1) < 15 and w > 20 and h > 20]

    print(f"        → 偵測到 {len(filtered)} 個前景圖案")
    return filtered, fg_bin

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 3a: 建立修復遮罩
# ─────────────────────────────────────────────────────────────────────────────

def create_inpaint_mask(image: Image.Image,
                        text_data: list,
                        fg_boxes: list,
                        fg_alpha: np.ndarray = None) -> Image.Image:
    """
    白色 = 需要修復（文字+前景物件），黑色 = 保留（背景）
    fg_alpha: rembg 的精確 alpha 遮罩，若提供則用於前景物件遮罩
    """
    w, h = image.size
    mask_arr = np.zeros((h, w), dtype=np.uint8)

    pad = INPAINT_PADDING

    # 文字區域（矩形遮罩）
    for td in text_data:
        x, y, bw, bh = td["rect"]
        x1 = max(0, x - pad);  y1 = max(0, y - pad)
        x2 = min(w, x + bw + pad); y2 = min(h, y + bh + pad)
        mask_arr[y1:y2, x1:x2] = 255

    # 前景物件遮罩：優先使用 rembg 精確 alpha，fallback 用邊界框
    if fg_alpha is not None:
        # 膨脹 alpha 遮罩以確保完整覆蓋
        k = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (pad * 2 + 1, pad * 2 + 1))
        dilated = cv2.dilate((fg_alpha > 30).astype(np.uint8) * 255, k, iterations=1)
        mask_arr = np.maximum(mask_arr, dilated)
    else:
        for x, y, bw, bh in fg_boxes:
            x1 = max(0, x - pad);  y1 = max(0, y - pad)
            x2 = min(w, x + bw + pad); y2 = min(h, y + bh + pad)
            mask_arr[y1:y2, x1:x2] = 255

    # 邊緣柔化 → 重新二值化
    blurred = cv2.GaussianBlur(mask_arr, (MASK_BLUR * 2 + 1, MASK_BLUR * 2 + 1), 0)
    final = (blurred > 25).astype(np.uint8) * 255
    return Image.fromarray(final)

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 3b: Stability AI Inpainting
# ─────────────────────────────────────────────────────────────────────────────

def _build_prompt(image: Image.Image, mask: Image.Image) -> str:
    """分析背景色，產生 inpainting 提示詞"""
    img_np  = np.array(image.convert("RGB"))
    mask_np = np.array(mask)
    bg_px   = img_np[mask_np < 128]
    if len(bg_px) == 0:
        return "clean smooth solid background, no text, no objects"

    r, g, b = bg_px.mean(axis=0).astype(int)
    if r > 180 and g > 180 and b > 180:
        tone = "light gray or white"
    elif r < 80 and g < 80 and b < 80:
        tone = "dark charcoal or black"
    elif r > g + 30 and r > b + 30:
        tone = "warm reddish toned"
    elif b > r + 30 and b > g + 30:
        tone = "cool blue toned"
    elif g > r + 20 and g > b + 20:
        tone = "green toned"
    else:
        tone = "neutral toned"

    return (f"seamless {tone} gradient background, "
            "smooth texture, no text, no icons, no objects, clean minimal")


def inpaint_with_stability(image: Image.Image,
                           mask: Image.Image,
                           api_key: str) -> Image.Image:
    """呼叫 Stability AI v2beta inpaint API"""
    # 解析度需為 64 的倍數
    ow, oh = image.size
    w = (ow // 64) * 64
    h = (oh // 64) * 64

    img_r  = image.resize((w, h), Image.LANCZOS)
    mask_r = mask.resize((w, h), Image.NEAREST)

    prompt = _build_prompt(image, mask)
    print(f"        → Inpaint 提示詞: \"{prompt}\"")

    img_buf, mask_buf = io.BytesIO(), io.BytesIO()
    img_r.save(img_buf,  format="PNG"); img_buf.seek(0)
    mask_r.save(mask_buf, format="PNG"); mask_buf.seek(0)

    resp = requests.post(
        "https://api.stability.ai/v2beta/stable-image/edit/inpaint",
        headers={"authorization": f"Bearer {api_key}", "accept": "image/*"},
        files={
            "image": ("image.png", img_buf,  "image/png"),
            "mask":  ("mask.png",  mask_buf, "image/png"),
        },
        data={
            "prompt":        prompt,
            "output_format": "png",
            "strength":      0.90,
            "seed":          42,
        },
        timeout=60,
    )

    if resp.status_code == 200:
        result = Image.open(io.BytesIO(resp.content)).convert("RGB")
        return result.resize((ow, oh), Image.LANCZOS)

    print(f"        [!] API 錯誤 {resp.status_code}: {resp.text[:200]}")
    return None


def fallback_inpaint(image: Image.Image, mask: Image.Image) -> Image.Image:
    """備用方案：OpenCV TELEA 演算法修復（不需 API）"""
    print("        → 使用 OpenCV 備用修復...")
    img_np  = np.array(image.convert("RGB"))
    mask_np = (np.array(mask.convert("L")) > 128).astype(np.uint8) * 255
    result  = cv2.inpaint(
        cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR),
        mask_np, inpaintRadius=18, flags=cv2.INPAINT_TELEA
    )
    return Image.fromarray(cv2.cvtColor(result, cv2.COLOR_BGR2RGB))


def regenerate_background(image: Image.Image,
                           mask: Image.Image,
                           api_key: str) -> Image.Image:
    """優先使用 Stability AI；失敗時退回 OpenCV"""
    if api_key:
        result = inpaint_with_stability(image, mask, api_key)
        if result is not None:
            return result
    return fallback_inpaint(image, mask)

# ─────────────────────────────────────────────────────────────────────────────
# 步驟 4: 組裝 PPTX
# ─────────────────────────────────────────────────────────────────────────────

def _estimate_font_pt(rect_h_px: int, img_h_px: int, slide_h_emu: int) -> int:
    """根據文字框高度估算字級（pt）"""
    h_emu = rect_h_px / img_h_px * slide_h_emu
    pt    = (h_emu / 12700) / 1.25   # 1 pt = 12700 EMU；行高約 1.25x
    return max(8, min(int(pt), 72))


def _sample_text_color(image: Image.Image, rect: tuple) -> tuple:
    """採樣文字框內像素，猜測文字顏色"""
    x, y, w, h = rect
    region = np.array(image.crop((x, y, x + w, y + h)).convert("RGB"))
    pixels = region.reshape(-1, 3)
    bright = pixels.mean(axis=1)

    # 判斷背景明暗
    median_bright = np.median(bright)
    if median_bright < 128:
        # 深色背景 → 文字可能是白色
        light = pixels[bright > 180]
        if len(light) > 5:
            c = light.mean(axis=0).astype(int)
            return (int(c[0]), int(c[1]), int(c[2]))
        return (255, 255, 255)
    else:
        # 淺色背景 → 文字可能是深色
        dark = pixels[bright < 100]
        if len(dark) > 5:
            c = dark.mean(axis=0).astype(int)
            return (int(c[0]), int(c[1]), int(c[2]))
        return (0, 0, 0)


def build_pptx(pages: list,
               text_per_page: list,
               fg_per_page: list,
               bg_images: list,
               output_path: str):

    prs = Presentation()

    ref_w, ref_h = pages[0].size
    SLIDE_W = 9144000                           # 10 英吋 EMU
    SLIDE_H = int(SLIDE_W * ref_h / ref_w)
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]         # 完全空白版面

    for page_idx, (page, text_data, fg_boxes, bg_img) in enumerate(
            zip(pages, text_per_page, fg_per_page, bg_images)):

        print(f"        組裝第 {page_idx + 1} 頁...")
        slide  = prs.slides.add_slide(blank_layout)
        img_w, img_h = page.size

        def ex(px): return px_scale(px, img_w, SLIDE_W)
        def ey(px): return px_scale(px, img_h, SLIDE_H)

        # ── 1. 背景圖（置底）──────────────────────────────────────────
        bg_path = ASSETS_DIR / f"bg_{page_idx:03d}.png"
        bg_img.save(bg_path)
        bg_shape = slide.shapes.add_picture(str(bg_path), 0, 0, SLIDE_W, SLIDE_H)
        # 移至最底層
        sp_tree = slide.shapes._spTree
        sp_tree.remove(bg_shape._element)
        sp_tree.insert(2, bg_shape._element)

        # ── 2. 前景圖案（可自由移動）─────────────────────────────────
        for obj_i, (ox, oy, ow, oh) in enumerate(fg_boxes):
            crop = page.crop((ox, oy, ox + ow, oy + oh))
            # 保留透明度（RGBA）
            crop_path = ASSETS_DIR / f"obj_{page_idx:03d}_{obj_i:02d}.png"
            crop.save(crop_path)
            slide.shapes.add_picture(
                str(crop_path), ex(ox), ey(oy), ex(ow), ey(oh)
            )

        # ── 3. 文字框（原始位置與字級）──────────────────────────────
        for td in text_data:
            x, y, w, h = td["rect"]
            left   = ex(x)
            top    = ey(y)
            width  = max(ex(w), ex(60))
            height = max(ey(h), ey(20))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf    = txBox.text_frame
            tf.word_wrap = False

            para = tf.paragraphs[0]
            run  = para.add_run()
            run.text = td["text"]

            font_pt = _estimate_font_pt(h, img_h, SLIDE_H)
            run.font.size = Pt(font_pt)

            r, g, b = _sample_text_color(page, td["rect"])
            run.font.color.rgb = RGBColor(r, g, b)

            # 移除文字框邊框與填色
            txBox.fill.background()

    prs.save(output_path)
    print(f"\n✓ 完成！輸出檔案：{output_path}")

# ─────────────────────────────────────────────────────────────────────────────
# 主程式
# ─────────────────────────────────────────────────────────────────────────────

def main():
    pdf_path    = sys.argv[1] if len(sys.argv) > 1 else PDF_PATH
    output_path = sys.argv[2] if len(sys.argv) > 2 else OUTPUT_PPTX

    if not Path(pdf_path).exists():
        print(f"錯誤：找不到檔案 → {pdf_path}")
        sys.exit(1)

    ASSETS_DIR.mkdir(exist_ok=True)

    api_key = STABILITY_API_KEY
    if not api_key:
        print("提示：未設定 STABILITY_API_KEY，將使用 OpenCV 備用修復（品質較低）")
        print("      設定方式：set STABILITY_API_KEY=your_key_here  (Windows)")

    # 步驟 0（可用 --pages N 只處理前 N 頁測試）
    max_pages = None
    if "--pages" in sys.argv:
        idx = sys.argv.index("--pages")
        max_pages = int(sys.argv[idx + 1])
    pages = pdf_to_images(pdf_path, last_page=max_pages)

    # 初始化 EasyOCR（只載入一次）
    print(f"\n[步驟 1] 初始化 OCR（語言：{OCR_LANGUAGES}，GPU：{USE_GPU}）...")
    reader = easyocr.Reader(OCR_LANGUAGES, gpu=USE_GPU)

    # 初始化 rembg session（只載入一次）
    print("[步驟 2] 初始化 rembg 前景偵測模型...")
    rembg_sess = rembg_session("u2net")

    text_per_page, fg_per_page, bg_images = [], [], []

    for i, page in enumerate(pages):
        print(f"\n══ 處理第 {i + 1}/{len(pages)} 頁 ══")
        page.save(ASSETS_DIR / f"original_{i:03d}.png")

        # 步驟 1
        print("[步驟 1] 偵測文字...")
        text_data = detect_text(page, reader)
        text_per_page.append(text_data)

        # 步驟 2
        print("[步驟 2] 偵測前景圖案（rembg）...")
        fg_boxes, fg_alpha = detect_foreground_objects(page, text_data, rembg_sess)
        fg_per_page.append(fg_boxes)
        # 儲存前景遮罩供除錯
        Image.fromarray(fg_alpha).save(ASSETS_DIR / f"fg_alpha_{i:03d}.png")

        # 步驟 3a
        print("[步驟 3] 建立修復遮罩...")
        mask = create_inpaint_mask(page, text_data, fg_boxes, fg_alpha)
        mask.save(ASSETS_DIR / f"mask_{i:03d}.png")

        # 步驟 3b（選項 B：暫時沿用原圖為背景，跳過 inpainting）
        print("[步驟 3] 使用原圖為背景（跳過 inpainting）...")
        bg = page.convert("RGB")
        bg.save(ASSETS_DIR / f"bg_{i:03d}.png")
        bg_images.append(bg)

        # API 限速保護
        if api_key and i < len(pages) - 1:
            time.sleep(1.5)

    # 步驟 4
    print("\n[步驟 4] 組裝 PPTX...")
    build_pptx(pages, text_per_page, fg_per_page, bg_images, output_path)


if __name__ == "__main__":
    main()
